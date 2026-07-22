from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 🎯 Create presentation
prs = Presentation()

# Utility function to add slide with title and content
def add_slide(title, content, highlight_words=[]):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Content
    content_placeholder = slide.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()

    for line in content.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(40, 40, 40)

        # Highlight important words
        for word in highlight_words:
            if word in line:
                run = p.add_run()
                run.text = f" ({word})"
                run.font.bold = True
                run.font.color.rgb = RGBColor(200, 0, 0)

# 🔹 Slides Content
slides_data = [
    ("SmartLock Project", 
     "A modern security solution with AI-powered Intruder Detection\nDeveloped as an IoT and Computer Vision project", 
     ["AI-powered", "Security"]),
    
    ("Problem Statement", 
     "• Passwords can be guessed or stolen\n• CCTV requires continuous monitoring\n• No instant alert system\n• Limited personal data protection", 
     ["Passwords", "CCTV", "alert"]),
    
    ("Our Solution – SmartLock", 
     "• Password-protected digital lock\n• Intruder detection via camera\n• Automatic photo capture on failed attempts\n• Email notification with evidence\n• Alarm trigger for security", 
     ["Intruder", "Automatic", "Alarm"]),
    
    ("Key Features", 
     "• Secure password system\n• Captures intruder selfie\n• Sends email alerts with image\n• Alarm sound on intrusion\n• Works offline after setup\n• Can be extended to mobile apps", 
     ["Secure", "intruder", "alerts", "Alarm"]),
    
    ("Real-Life Applications", 
     "• Home door security\n• Office and lab access\n• School computer labs\n• Personal devices\n• Future: Mobile app protection for WhatsApp, Instagram", 
     ["Home", "Office", "Mobile"]),
    
    ("Limitations", 
     "• Requires webcam support\n• Needs internet for email alerts\n• Currently limited to PC/laptop\n• Play Store release needs license", 
     ["internet", "PC/laptop", "Play Store"]),
    
    ("Conclusion", 
     "SmartLock is a step towards modern digital security\n✔ Cost-effective\n✔ Easy to use\n✔ Provides real-time alerts\n✔ Can evolve into mobile app in future", 
     ["security", "real-time", "mobile"])
]

# Add all slides
for title, content, highlights in slides_data:
    add_slide(title, content, highlights)

# Save presentation
prs.save("SmartLock_Project_Presentation.pptx")
print("✅ Presentation created: SmartLock_Project_Presentation.pptx")
